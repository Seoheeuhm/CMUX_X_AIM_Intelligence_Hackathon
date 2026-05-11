import io
import json
import os
import re
import tempfile
import uuid
from datetime import datetime
from typing import Any, Dict, List

import anthropic
import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from fastapi import Depends, FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from starlette.background import BackgroundTask

from auth import get_current_user, get_optional_user
from db import (
    payment_confirm,
    payment_create,
    payment_fail,
    payment_get_status,

    profile_can_generate,
    profile_get,
    profile_increment_gen,
    profile_upgrade_to_pro,
    session_get,
    session_set,
)
from payment import toss_confirm

load_dotenv()
print("KEY LOADED:", bool(os.getenv("ANTHROPIC_API_KEY")))

try:
    from pptx import Presentation
    from pptx.util import Cm, Pt
    from pptx.dml.color import RGBColor
    PPTX_OK = True
except ImportError:
    PPTX_OK = False
    RGBColor = Cm = Pt = None

try:
    import PyPDF2
    PDF_OK = True
except ImportError:
    PDF_OK = False

app = FastAPI(title="AI Portfolio Generator", docs_url=None, redoc_url=None)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

portfolio_history: List[Dict[str, Any]] = []

SARAMIN_DUMMY = [
    {
        "id": "d001",
        "title": "백엔드 개발자 (Python/FastAPI)",
        "company": "테크스타트업 A",
        "deadline": "2025-05-31",
        "job_type": "정규직",
        "experience": "경력 3-5년",
        "location": "서울 강남구",
        "description": "• Python/FastAPI 기반 RESTful API 개발\n• PostgreSQL, Redis 운영 경험\n• AWS EC2/RDS/S3 사용 경험\n• CI/CD(Github Actions) 경험 우대\n• 스타트업 경험 우대",
    },
    {
        "id": "d002",
        "title": "프론트엔드 개발자 (React)",
        "company": "이커머스 플랫폼 B",
        "deadline": "2025-06-15",
        "job_type": "정규직",
        "experience": "경력 2-4년",
        "location": "서울 마포구",
        "description": "• React/TypeScript 기반 웹 애플리케이션 개발\n• Next.js 경험 우대\n• 반응형 UI 구현 경험\n• 성능 최적화 경험 보유자 우대\n• 디자인 시스템 구축 경험",
    },
    {
        "id": "d003",
        "title": "서비스 기획자 (PM)",
        "company": "핀테크 스타트업 C",
        "deadline": "2025-05-25",
        "job_type": "정규직",
        "experience": "경력 3년 이상",
        "location": "서울 여의도",
        "description": "• 모바일 앱/웹 서비스 기획 및 운영\n• 사용자 리서치 및 데이터 기반 의사결정\n• 개발팀과의 협업 경험\n• Figma/Notion 활용 능숙\n• 금융/결제 도메인 경험 우대",
    },
    {
        "id": "d004",
        "title": "UX/UI 디자이너",
        "company": "SaaS 기업 D",
        "deadline": "2025-06-30",
        "job_type": "정규직",
        "experience": "경력 2-5년",
        "location": "서울 성수동",
        "description": "• Figma 활용 UI 디자인 및 프로토타이핑\n• 사용자 여정 지도 및 와이어프레임 제작\n• 디자인 시스템 구축 및 관리\n• A/B 테스트 설계 및 분석 경험\n• B2B SaaS 경험 우대",
    },
    {
        "id": "d005",
        "title": "퍼포먼스 마케터",
        "company": "이커머스 스타트업 E",
        "deadline": "2025-06-10",
        "job_type": "정규직",
        "experience": "경력 2-4년",
        "location": "서울 강남구",
        "description": "• Meta/Google/카카오 광고 운영 경험\n• 데이터 분석 (GA4, Amplitude) 능숙\n• ROAS, CAC, LTV 지표 최적화 경험\n• SQL 기초 쿼리 가능\n• 그로스해킹 마인드셋 보유자",
    },
    {
        "id": "d006",
        "title": "데이터 분석가",
        "company": "헬스케어 스타트업 F",
        "deadline": "2025-07-01",
        "job_type": "정규직",
        "experience": "경력 2년 이상",
        "location": "서울 종로구",
        "description": "• Python/SQL 기반 데이터 분석\n• Tableau/Looker Studio 시각화\n• A/B 테스트 설계 및 통계 분석\n• 머신러닝 기초 지식 우대\n• 헬스케어 도메인 경험 우대",
    },
    {
        "id": "d007",
        "title": "iOS 개발자 (Swift)",
        "company": "모빌리티 서비스 G",
        "deadline": "2025-06-20",
        "job_type": "정규직",
        "experience": "경력 3년 이상",
        "location": "서울 서초구",
        "description": "• Swift/SwiftUI 기반 iOS 앱 개발\n• REST API 연동 및 데이터 처리\n• 앱 성능 최적화 경험\n• TDD/BDD 경험 우대\n• 지도/위치 서비스 개발 경험 우대",
    },
    {
        "id": "d008",
        "title": "콘텐츠 마케터",
        "company": "에듀테크 기업 H",
        "deadline": "2025-05-28",
        "job_type": "정규직",
        "experience": "경력 1-3년",
        "location": "서울 강남구",
        "description": "• 블로그/SNS 채널 콘텐츠 기획 및 제작\n• SEO 최적화 전략 수립\n• 이메일 마케팅 캠페인 운영\n• CRM 툴 활용 경험\n• 교육 분야 콘텐츠 경험 우대",
    },
]


def client() -> anthropic.Anthropic:
    key = os.getenv("ANTHROPIC_API_KEY")
    if not key:
        raise HTTPException(500, "ANTHROPIC_API_KEY 환경변수가 설정되지 않았습니다.")
    return anthropic.Anthropic(api_key=key)

def extract_pptx(data: bytes) -> str:
    if not PPTX_OK:
        raise HTTPException(400, "python-pptx가 설치되지 않았습니다.")
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_pdf(data: bytes) -> str:
    if not PDF_OK:
        raise HTTPException(400, "PyPDF2가 설치되지 않았습니다.")
    reader = PyPDF2.PdfReader(io.BytesIO(data))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            parts.append(f"[페이지 {i}]\n{text.strip()}")
    return "\n\n".join(parts)


def scrape(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; PortfolioBot/1.0)"}
        r = requests.get(url, headers=headers, timeout=15)
        soup = BeautifulSoup(r.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "meta", "link"]):
            tag.decompose()
        lines = [ln.strip() for ln in soup.get_text("\n").splitlines() if ln.strip()]
        return "\n".join(lines[:700])
    except Exception:
        return ""


def clean_html(raw: str) -> str:
    raw = raw.strip()
    raw = re.sub(r"^```(?:html)?\s*", "", raw, flags=re.IGNORECASE)
    raw = re.sub(r"\s*```$", "", raw)
    return raw.strip()


def call_claude(c: anthropic.Anthropic, prompt: str, max_tokens: int = 8192, max_rounds: int = 4) -> str:
    messages = [{"role": "user", "content": prompt}]
    accumulated = ""

    for _ in range(max_rounds):
        msg = c.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=max_tokens,
            messages=messages,
        )
        chunk = msg.content[0].text
        accumulated += chunk

        if msg.stop_reason != "max_tokens":
            break

        messages.append({"role": "assistant", "content": chunk})
        messages.append({
            "role": "user",
            "content": (
                "HTML이 중간에 잘렸습니다. "
                "이미 작성한 내용은 반복하지 말고, "
                "잘린 부분부터 이어서 나머지 HTML을 완성해주세요. "
                "모든 열린 태그를 올바르게 닫아주세요."
            ),
        })

    return clean_html(accumulated)

def parse_materials(text: str, source: str) -> dict:
    c = client()
    msg = c.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=2048,
        messages=[{"role": "user", "content": (
            f"다음 포트폴리오 텍스트를 분석해서 JSON만 반환하세요 (코드블록·설명 없이).\n"
            f"출처: {source}\n\n텍스트:\n{text[:8000]}\n\n"
            '반환 형식:\n{"projects":[{"name":"","role":"","tech_stack":[],"achievements":[],"description":""}],'
            '"skills":[],"summary":""}'
        )}],
    )
    raw = msg.content[0].text.strip()
    m = re.search(r"\{[\s\S]*\}", raw)
    if m:
        try:
            return json.loads(m.group())
        except Exception:
            pass
    return {"projects": [], "skills": [], "summary": text[:300]}

# ── Prompt builders ───────────────────────────────────────────────────────────

TEMPLATES = {
    "developer": {
        "s1_tag": "기술스택 중심 (언어·프레임워크·인프라)",
        "s3_title": "기술 역량 요약",
        "s4_metric": "성능 개선율·코드 품질·배포 횟수 등 기술 지표",
        "s5_title": "기술적 문제해결 역량",
    },
    "planner": {
        "s1_tag": "기획 방법론 중심 (리서치·UX·지표 설계)",
        "s3_title": "기획 역량 요약",
        "s4_metric": "MAU·전환율·NPS·기획 범위 등 서비스 지표",
        "s5_title": "서비스 기획 역량",
    },
    "designer": {
        "s1_tag": "툴·결과물 중심 (Figma·브랜딩·UX 개선)",
        "s3_title": "디자인 역량 요약",
        "s4_metric": "사용성 개선·브랜드 임팩트·결과물 수 등 시각 지표",
        "s5_title": "디자인 프로세스 역량",
    },
    "marketer": {
        "s1_tag": "채널·지표 중심 (캠페인·ROAS·CTR·세그먼트)",
        "s3_title": "채널 경험 요약",
        "s4_metric": "CTR·ROAS·CAC·세그먼트 수 등 퍼포먼스 지표",
        "s5_title": "데이터 기반 전략 역량",
    },
}

COMMON_RULES = """
공통 HTML 규칙:
- 섹션: <section class="section">
- 섹션 제목: <h2 class="section-title">제목</h2>
- 항목 카드: <div class="item">
- 기술 태그: <span class="tech-tag">태그</span>
- 채용공고 키워드 강조(Track B만): <strong class="kw">키워드</strong>
- 마크다운·코드블록 금지, 순수 HTML만 반환
- <div class="portfolio-content"> 래퍼 없이 섹션만 반환

[이력서 디자인 원칙]
- 배경: 흰색, 텍스트: 주로 #1a1a1a(거의 검정)
- 중요 키워드: <strong class="kw"> 또는 <span class="highlight-box">태그</span> 사용
- 수치·성과: 반드시 <strong class="kw">로 강조 (예: <strong class="kw">MAU 3만</strong>)
- 섹션 제목: 대문자 영문 병기 권장 (예: 주요 프로젝트 PROJECTS)
- 본문: 14px, 줄간격 넉넉하게, 읽기 쉽게
- 태그 칩(.tech-tag): 회색 배경, 검정 텍스트로 눈에 띄게
- 절대 다크 배경(#1e1e1e, #2D2D2D 등) 사용 금지
"""


def build_prompt_track_a(mat_json: str, interview_txt: str, tmpl: dict) -> str:
    return f"""당신은 이력서·포트폴리오 작성 전문가입니다.
아래 지원자 정보를 바탕으로 깔끔한 흰색 배경 이력서 HTML을 생성하세요.
중요한 수치·키워드는 <strong class='kw'>로 강조하고, 기술명은 <span class='tech-tag'>로 표시하세요.
절대 다크 배경을 사용하지 마세요.

[지원자 정보]
파싱된 포트폴리오: {mat_json}
인터뷰 답변: {interview_txt or "(없음)"}

[양식]
{tmpl['s1_tag']}

{COMMON_RULES}

아래 3개 섹션만 생성하세요 (순서 유지):

SECTION 1 — 자기소개 & 핵심 역량
  <section class="section s1">: 강점 3가지를 자연스러운 소개 문장으로. 수치는 <strong class="kw">로 강조.

SECTION 2 — 주요 프로젝트 PROJECTS
  <section class="section s2">: 프로젝트별 카드. 각 카드에 역할·기술스택·성과 포함.
  성과 수치는 반드시 <strong class="kw">로 강조. 기술명은 <span class="tech-tag">로 표시.

SECTION 3 — {tmpl['s3_title']}
  <section class="section s3">: 보유 기술/역량을 카테고리별로. tech-tag 적극 활용.
"""


def build_prompt_track_b(mat_json: str, interview_txt: str,
                          tmpl: dict, job_title: str, job_posting: str) -> str:
    return f"""당신은 채용 전문가 겸 이력서·포트폴리오 작성 전문가입니다.
아래 지원자 정보와 채용공고를 바탕으로 깔끔한 흰색 배경 맞춤형 이력서 HTML을 생성하세요.
채용공고 핵심 키워드는 <strong class='kw'>로 강조하고, 기술명은 <span class='tech-tag'>로 표시하세요.
절대 다크 배경을 사용하지 마세요.

[지원자 정보]
파싱된 포트폴리오: {mat_json}
인터뷰 답변: {interview_txt or "(없음)"}

[채용 정보]
지원 직무: {job_title}
채용공고: {job_posting[:2500]}

[양식]
{tmpl['s1_tag']}

{COMMON_RULES}

아래 5개 섹션만 생성하세요 (순서 유지):

SECTION 1 — 지원자 프로필 & 직무 적합성
  <section class="section s1">: 채용공고 요구사항에 맞춘 자기소개 + 핵심 강점 3가지.
  채용공고 키워드를 <strong class="kw">로 강조.

SECTION 2 — 주요 프로젝트 PROJECTS
  <section class="section s2">: 프로젝트별 카드. 채용공고 관련 내용 우선 배치.
  성과 수치는 <strong class="kw">로 강조. 채용공고 키워드를 <strong class="kw">로 강조.

SECTION 3 — {tmpl['s3_title']}
  <section class="section s3">: 보유 기술/역량 카테고리별 정리. tech-tag 활용.

SECTION 4 — 채용공고 요구 역량 매핑
  <section class="section s4">: 채용공고 요구 역량과 내 경험을 1:1 매핑 카드로.
  형식: [요구역량] → [나의 경험 근거]

SECTION 5 — {tmpl['s5_title']}
  <section class="section s5">: 대표적인 문제→해결→결과 사례 2~3개.
  {tmpl['s4_metric']} 포함.
"""


def build_prompt_second_call(mat_json: str, interview_txt: str, retro_txt: str,
                              has_jd: bool, job_title: str = "", job_posting: str = "") -> str:
    jd_block = ""
    if has_jd:
        jd_block = f"[채용공고] 지원직무: {job_title}\n{job_posting[:1500]}"

    if has_jd:
        sections = f"""
SECTION 6 — 회고 & 성장 스토리
  <section class="section s6">: 프로젝트별 회고를 자연스러운 문장으로.
  1) 핵심 인사이트  2) 파악된 업무 스타일  3) 다음 프로젝트 적용 계획.

SECTION 7 — 핵심 성과 지표 요약
  <section class="section s7">: <table> 태그로 5~8행 구성.
  컬럼: 항목 | 수치 | 직무 연계 포인트.
  테이블 헤더 배경 #1a1a1a, 흰색 텍스트, 셀 텍스트 #1a1a1a.

SECTION 8 — 미보유 역량 & 학습 계획
  <section class="section s8">: 채용공고 기준 미보유 역량 2~3개와 구체적 보완 방법.
"""
    else:
        sections = """
SECTION 4 — 회고 & 성장 스토리
  <section class="section s4">: 프로젝트별 회고를 자연스러운 문장으로.
  1) 핵심 인사이트  2) 파악된 업무 스타일  3) 다음 프로젝트 적용 계획.
"""

    return f"""흰색 배경 이력서 스타일 유지. 다크 배경 사용 금지.
수치·성과는 <strong class='kw'>로 강조.
아래 내용을 바탕으로 추가 섹션만 HTML로 생성하세요.
순수 HTML만 반환, 코드블록 금지.

[지원자 정보]
{mat_json}
인터뷰 답변: {interview_txt or "(없음)"}
회고: {retro_txt or "(없음)"}
{jd_block}

{COMMON_RULES}

아래 섹션만 생성하세요 (다른 내용 절대 추가 금지):
{sections}
"""


# ── Pydantic models ───────────────────────────────────────────────────────────

class UrlReq(BaseModel):
    url: str

class TextReq(BaseModel):
    text: str

class InterviewStartReq(BaseModel):
    session_id: str

class InterviewAnswerReq(BaseModel):
    session_id: str
    answer: str

class RetroReq(BaseModel):
    session_id: str
    retro_type: str
    retro_data: list

class GenerateReq(BaseModel):
    session_id: str
    portfolio_type: str
    job_title: str = ""
    job_posting: str = ""

class PaymentCreateReq(BaseModel):
    amount: int = 9900

class PaymentConfirmReq(BaseModel):
    payment_key: str
    order_id: str
    amount: int


class SlideCard(BaseModel):
    title: str = ""
    body: str = ""


class Slide(BaseModel):
    id: str
    layout: str               # title | section | content | project | closing
    title: str = ""
    subtitle: str = ""
    content: str = ""
    bullets: list = []
    cards: list = []
    role: str = ""
    description: str = ""
    tech_stack: list = []
    achievements: list = []


class PptPreviewReq(BaseModel):
    session_id: str
    portfolio_type: str       # developer | planner | designer | marketer
    job_title: str = ""
    job_posting: str = ""


class PptDownloadReq(BaseModel):
    session_id: str
    portfolio_type: str
    slides: list


# ── Static files ──────────────────────────────────────────────────────────────

app.mount("/static", StaticFiles(directory="static"), name="static")
app.mount("/image", StaticFiles(directory="image"), name="image")

_NO_CACHE = {"Cache-Control": "no-cache, no-store, must-revalidate", "Pragma": "no-cache"}

@app.get("/")
def root():
    return FileResponse("static/landing.html", headers=_NO_CACHE)

@app.get("/app")
def app_page():
    return FileResponse("static/index.html", headers=_NO_CACHE)

@app.get("/generator")
def generator_page():
    return FileResponse("static/index.html", headers=_NO_CACHE)

@app.get("/docs")
def docs_page():
    return FileResponse("static/docs.html", headers=_NO_CACHE)

@app.get("/qna")
def qna_page():
    return FileResponse("static/qna.html", headers=_NO_CACHE)


# ── Auth ──────────────────────────────────────────────────────────────────────

@app.get("/auth/login")
def auth_login():
    supabase_url = os.environ["SUPABASE_URL"]
    app_url = os.environ.get("APP_URL", "http://localhost:8000")
    redirect_to = f"{app_url}/auth/callback"
    oauth_url = (
        f"{supabase_url}/auth/v1/authorize"
        f"?provider=google"
        f"&redirect_to={redirect_to}"
    )
    return RedirectResponse(url=oauth_url)


@app.get("/auth/callback")
def auth_callback():
    return FileResponse("static/auth-callback.html")


@app.get("/auth/me")
def auth_me(user: dict = Depends(get_current_user)):
    profile = profile_get(user["sub"])
    return {
        "id": user["sub"],
        "email": user.get("email", ""),
        "profile": profile,
    }


@app.get("/config")
def get_config():
    return {"toss_client_key": os.environ.get("TOSS_CLIENT_KEY", "")}


# ── Job listings ──────────────────────────────────────────────────────────────

@app.get("/api/jobs/seoul")
def get_seoul_jobs(page: int = 1, per_page: int = 30, keyword: str = ""):
    key = os.getenv("SEOUL_API_KEY", "")
    if not key:
        return {"jobs": [], "total": 0, "page": page, "per_page": per_page,
                "error": "SEOUL_API_KEY 환경변수가 설정되지 않았습니다."}
    start = (page - 1) * per_page + 1
    end = page * per_page
    url = f"http://openapi.seoul.go.kr:8088/{key}/json/GetJobInfo/{start}/{end}/"
    try:
        r = requests.get(url, timeout=10)
        data = r.json()
        if "RESULT" in data and "GetJobInfo" not in data:
            msg = data["RESULT"].get("MESSAGE", "데이터를 불러올 수 없습니다.")
            return {"jobs": [], "total": 0, "page": page, "per_page": per_page, "error": msg}
        info = data.get("GetJobInfo", {})
        result_code = info.get("RESULT", {}).get("CODE", "")
        if result_code and result_code != "INFO-000":
            msg = info.get("RESULT", {}).get("MESSAGE", "데이터 없음")
            return {"jobs": [], "total": 0, "page": page, "per_page": per_page, "error": msg}
        rows = info.get("row", [])
        total = int(info.get("list_total_count", len(rows)))
        return {"jobs": rows, "total": total, "page": page, "per_page": per_page}
    except Exception as e:
        return {"jobs": [], "total": 0, "page": page, "per_page": per_page,
                "error": f"API 호출 실패: {str(e)}"}


@app.get("/api/jobs/saramin")
def get_saramin_jobs():
    key = os.getenv("SARAMIN_API_KEY", "")
    if not key:
        return {"jobs": SARAMIN_DUMMY, "is_dummy": True}
    url = "https://oapi.saramin.co.kr/job-search"
    params = {"access-key": key, "count": 20, "start": 1, "fields": "all"}
    try:
        r = requests.get(url, params=params, timeout=10,
                         headers={"Accept": "application/json"})
        data = r.json()
        raw_jobs = data.get("jobs", {}).get("job", [])
        formatted = []
        for job in raw_jobs:
            pos = job.get("position", {})
            formatted.append({
                "id": str(job.get("id", "")),
                "title": pos.get("title", ""),
                "company": job.get("company", {}).get("detail", {}).get("name", ""),
                "deadline": job.get("expiration-date", ""),
                "job_type": pos.get("job-type", {}).get("name", ""),
                "experience": pos.get("experience-level", {}).get("name", ""),
                "location": pos.get("location", {}).get("name", ""),
                "description": pos.get("title", "") + "\n" + pos.get("required-education-level", {}).get("name", ""),
            })
        if not formatted:
            return {"jobs": SARAMIN_DUMMY, "is_dummy": True}
        return {"jobs": formatted, "is_dummy": False}
    except Exception:
        return {"jobs": SARAMIN_DUMMY, "is_dummy": True}


# ── Portfolio history ─────────────────────────────────────────────────────────

@app.get("/api/history")
def get_history():
    return {"history": portfolio_history[:3]}


# ── Upload ────────────────────────────────────────────────────────────────────

@app.post("/upload/file")
async def upload_file(files: List[UploadFile] = File(...)):
    sid = str(uuid.uuid4())
    materials = []
    for f in files:
        name = f.filename or ""
        data = await f.read()
        ext = name.lower().rsplit(".", 1)[-1] if "." in name else ""
        if ext in ("pptx", "ppt"):
            raw = extract_pptx(data)
        elif ext == "pdf":
            raw = extract_pdf(data)
        else:
            raise HTTPException(400, f"지원하지 않는 형식: {name}")
        if not raw.strip():
            raise HTTPException(400, f"텍스트를 추출할 수 없습니다: {name}")
        parsed = parse_materials(raw, name)
        materials.append({"source": name, "raw": raw[:5000], "parsed": parsed})

    session_set(sid, {
        "materials": materials,
        "retro": {},
        "questions": [],
        "answers": [],
        "q_index": 0,
    })
    return {"session_id": sid, "parsed": [m["parsed"] for m in materials]}


@app.post("/upload/notion")
def upload_notion(req: UrlReq):
    raw = scrape(req.url)
    if not raw:
        raise HTTPException(400, "노션 페이지를 가져올 수 없습니다. 공개 링크인지 확인해주세요.")
    sid = str(uuid.uuid4())
    parsed = parse_materials(raw, "Notion")
    session_set(sid, {
        "materials": [{"source": "Notion", "raw": raw[:5000], "parsed": parsed}],
        "retro": {},
        "questions": [],
        "answers": [],
        "q_index": 0,
    })
    return {"session_id": sid, "parsed": [parsed]}


@app.post("/upload/text")
def upload_text(req: TextReq):
    if not req.text.strip():
        raise HTTPException(400, "텍스트를 입력해주세요.")
    sid = str(uuid.uuid4())
    parsed = parse_materials(req.text, "직접 입력")
    session_set(sid, {
        "materials": [{"source": "직접 입력", "raw": req.text[:5000], "parsed": parsed}],
        "retro": {},
        "questions": [],
        "answers": [],
        "q_index": 0,
    })
    return {"session_id": sid, "parsed": [parsed]}


# ── Retrospective ─────────────────────────────────────────────────────────────

@app.post("/retro/save")
def retro_save(req: RetroReq):
    sess = session_get(req.session_id)
    if sess is None:
        raise HTTPException(404, "세션을 찾을 수 없습니다.")
    sess["retro"] = {"type": req.retro_type, "data": req.retro_data}
    session_set(req.session_id, sess)
    return {"ok": True}


# ── Interview ─────────────────────────────────────────────────────────────────

@app.post("/interview/start")
def interview_start(req: InterviewStartReq):
    sess = session_get(req.session_id)
    if sess is None:
        raise HTTPException(404, "세션을 찾을 수 없습니다.")
    summary = json.dumps([m["parsed"] for m in sess["materials"]], ensure_ascii=False)

    c = client()
    msg = c.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=600,
        messages=[{"role": "user", "content": (
            f"포트폴리오 내용:\n{summary}\n\n"
            "이 경험을 더 깊이 파악하는 꼬리 질문 3개를 JSON만 반환하세요 (설명 없이):\n"
            '{"questions":["질문1","질문2","질문3"]}\n\n'
            "질문 주제: 1) 가장 어려운 기술적/업무적 문제와 해결법 "
            "2) 구체적 성과 수치 또는 임팩트 3) 팀 내 핵심 기여와 역할"
        )}],
    )
    raw = msg.content[0].text.strip()
    m = re.search(r"\{[\s\S]*\}", raw)
    qs = None
    if m:
        try:
            qs = json.loads(m.group()).get("questions")
        except Exception:
            pass
    if not qs or len(qs) < 3:
        qs = [
            "가장 큰 기술적 챌린지는 무엇이었고 어떻게 해결하셨나요?",
            "가장 자랑스러운 성과의 구체적인 수치나 임팩트가 있다면 알려주세요.",
            "팀 프로젝트에서 본인의 핵심 기여와 역할은 무엇이었나요?",
        ]
    sess["questions"] = qs
    sess["q_index"] = 0
    session_set(req.session_id, sess)
    return {"question": qs[0], "number": 1, "total": len(qs)}


@app.post("/interview/answer")
def interview_answer(req: InterviewAnswerReq):
    sess = session_get(req.session_id)
    if sess is None:
        raise HTTPException(404, "세션을 찾을 수 없습니다.")
    sess["answers"] = sess.get("answers", []) + [req.answer]
    sess["q_index"] = sess.get("q_index", 0) + 1
    session_set(req.session_id, sess)
    if sess["q_index"] >= len(sess["questions"]):
        return {"done": True}
    q = sess["questions"][sess["q_index"]]
    return {"done": False, "question": q, "number": sess["q_index"] + 1, "total": len(sess["questions"])}


# ── Scrape job posting ────────────────────────────────────────────────────────

@app.post("/scrape")
def scrape_url(req: UrlReq):
    text = scrape(req.url)
    if not text:
        raise HTTPException(400, "URL을 파싱할 수 없습니다.")
    return {"text": text}


# ── Generate portfolio ────────────────────────────────────────────────────────

@app.post("/generate")
def generate(req: GenerateReq, user: dict = Depends(get_current_user)):
    user_id = user["sub"]

    if not profile_can_generate(user_id):
        raise HTTPException(
            status_code=402,
            detail="생성 횟수를 모두 사용했습니다. Pro 플랜으로 업그레이드하세요.",
        )

    sess = session_get(req.session_id)
    if sess is None:
        raise HTTPException(404, "세션을 찾을 수 없습니다.")

    mat_json = json.dumps([m["parsed"] for m in sess["materials"]], ensure_ascii=False)

    interview_txt = ""
    for i, (q, a) in enumerate(zip(sess.get("questions", []), sess.get("answers", [])), 1):
        interview_txt += f"Q{i}. {q}\nA. {a}\n\n"

    retro_txt = ""
    if sess.get("retro"):
        r = sess["retro"]
        rtype = r.get("type", "")
        rdata = r.get("data", [])
        if rtype and rtype != "skip" and rdata:
            retro_txt = f"회고 유형: {rtype}\n"
            if isinstance(rdata, list):
                for pr in rdata:
                    pname = pr.get("project", "프로젝트")
                    retro_txt += f"\n[프로젝트: {pname}]\n"
                    for k, val in pr.items():
                        if k != "project" and val:
                            retro_txt += f"  {k}: {val}\n"
            elif isinstance(rdata, dict):
                for k, val in rdata.items():
                    if val:
                        retro_txt += f"  {k}: {val}\n"

    tmpl = TEMPLATES.get(req.portfolio_type, TEMPLATES["developer"])
    has_jd = bool(req.job_title and req.job_posting and req.job_posting.strip())

    if has_jd:
        prompt1 = build_prompt_track_b(
            mat_json, interview_txt, tmpl, req.job_title, req.job_posting
        )
    else:
        prompt1 = build_prompt_track_a(mat_json, interview_txt, tmpl)

    c = client()
    html1 = call_claude(c, prompt1)

    prompt2 = build_prompt_second_call(
        mat_json, interview_txt, retro_txt,
        has_jd, req.job_title, req.job_posting,
    )
    html2 = call_claude(c, prompt2)

    combined = html1 + "\n" + html2
    if '<div class="portfolio-content">' not in combined:
        combined = f'<div class="portfolio-content">{combined}</div>'

    type_names = {"developer": "개발자", "planner": "기획자", "designer": "디자이너", "marketer": "마케터"}
    portfolio_history.insert(0, {
        "id": str(uuid.uuid4()),
        "created_at": datetime.now().strftime("%Y.%m.%d %H:%M"),
        "job_title": req.job_title if req.job_title else "경험 정제 카드",
        "portfolio_type": req.portfolio_type,
        "portfolio_type_name": type_names.get(req.portfolio_type, req.portfolio_type),
        "track": "B" if has_jd else "A",
        "full_html": combined,
    })
    if len(portfolio_history) > 50:
        portfolio_history[:] = portfolio_history[:50]
    profile_increment_gen(user_id)

    return {
        "portfolio_html": combined,
        "job_title": req.job_title,
        "portfolio_type": req.portfolio_type,
        "track": "B" if has_jd else "A",
    }


# ── PPT Design System ─────────────────────────────────────────────────────────

_PPT_BG       = "1e1e1e"
_PPT_BG_CARD  = "2a2a2a"
_PPT_WHITE    = "FFFFFF"
_PPT_GRAY_MED = "999999"
_PPT_GRAY_DIM = "888888"
_PPT_GRAY_DARK= "666666"

_PPT_ACCENTS = {
    "developer": "6366F1",
    "planner":   "6366F1",
    "designer":  "6366F1",
    "marketer":  "6366F1",
}


def _hex_rgb(s: str):
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _ppt_rect(slide, x: float, y: float, w: float, h: float, fill: str):
    shp = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_rgb(fill)
    shp.line.fill.background()
    return shp


def _ppt_txt(slide, text: str, x: float, y: float, w: float, h: float,
             size: float, bold: bool = False, color: str = "FFFFFF"):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = str(text)[:300]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = _hex_rgb(color)
    return tb


def _ppt_dark_slide(slide, sd: dict, layout: str, accent: str):
    W, H = 33.87, 19.05
    _ppt_rect(slide, 0, 0, W, H, _PPT_BG)
    _ppt_rect(slide, 0.6, 2.5, 0.4, 6.0, _PPT_WHITE)
    _ppt_txt(slide, sd.get("title", ""), 1.8, 3.0, W - 3.0, 3.5, 34, bold=True, color=_PPT_WHITE)
    sub = sd.get("subtitle") or sd.get("content") or ""
    if sub:
        _ppt_txt(slide, sub, 1.8, 7.0, W - 3.0, 2.0, 16, color=_PPT_GRAY_MED)


def _ppt_project_slide(slide, sd: dict, accent: str):
    W, H = 33.87, 19.05
    _ppt_rect(slide, 0, 0, W, H, _PPT_BG)
    _ppt_rect(slide, 0, 0, W * 0.4, H, _PPT_BG_CARD)
    _ppt_txt(slide, sd.get("title", ""), 1.0, 1.2, W * 0.4 - 2.0, 2.5, 22, bold=True, color=_PPT_WHITE)
    date_str = sd.get("date", "")
    if date_str:
        _ppt_txt(slide, date_str, 1.0, 3.9, W * 0.4 - 2.0, 0.8, 10, color=_PPT_GRAY_DIM)
    role = sd.get("role", "")
    if role:
        _ppt_txt(slide, role, 1.0, 4.8, W * 0.4 - 2.0, 0.8, 11, color=_PPT_GRAY_MED)

    right_x = W * 0.4 + 0.8
    right_w = W - right_x - 1.0
    ry = 1.2
    desc = sd.get("description", "")
    if desc:
        _ppt_txt(slide, desc[:220], right_x, ry, right_w, 3.5, 11, color=_PPT_GRAY_MED)
        ry += 3.8
    tech = sd.get("tech_stack", [])
    ach  = sd.get("achievements", [])
    if tech:
        _ppt_txt(slide, "STACK", right_x, ry, right_w, 0.7, 9, bold=True, color=_PPT_GRAY_DARK)
        ry += 0.8
        _ppt_txt(slide, "  ·  ".join(str(t) for t in tech[:8]), right_x, ry, right_w, 1.4, 10, color=_PPT_GRAY_MED)
        ry += 1.8
    if ach:
        _ppt_txt(slide, "ACHIEVEMENTS", right_x, ry, right_w, 0.7, 9, bold=True, color=_PPT_GRAY_DARK)
        ry += 0.8
        for a in ach[:4]:
            if ry > H - 1.5:
                break
            _ppt_txt(slide, f"— {str(a)[:80]}", right_x, ry, right_w, 0.9, 10, color=_PPT_WHITE)
            ry += 1.0


def _ppt_content_slide(slide, sd: dict, accent: str):
    W, H = 33.87, 19.05
    _ppt_rect(slide, 0, 0, W, H, _PPT_BG)
    _ppt_txt(slide, sd.get("title", "").upper(), 1.2, 0.8, W - 2.0, 1.4, 13, bold=True, color=_PPT_WHITE)
    _ppt_rect(slide, 1.2, 2.4, W - 2.4, 0.05, _PPT_GRAY_DARK)

    cards   = sd.get("cards", [])
    bullets = sd.get("bullets", [])

    if cards:
        n = min(len(cards), 3)
        card_w = (W - 2.0 - (n - 1) * 0.5) / n
        for i, card in enumerate(cards[:3]):
            cx = 1.0 + i * (card_w + 0.5)
            cy, ch = 3.0, H - 4.2
            _ppt_rect(slide, cx, cy, card_w, ch, _PPT_BG_CARD)
            c = card if isinstance(card, dict) else {"title": str(card), "body": ""}
            _ppt_txt(slide, c.get("title", ""), cx + 0.4, cy + 0.5, card_w - 0.8, 1.0, 13, bold=True, color=_PPT_WHITE)
            if c.get("body"):
                _ppt_txt(slide, c["body"], cx + 0.4, cy + 1.7, card_w - 0.8, ch - 2.0, 11, color=_PPT_GRAY_MED)
    elif bullets:
        by = 2.8
        for b in bullets[:7]:
            if by > H - 1.5:
                break
            _ppt_rect(slide, 1.2, by + 0.35, 0.25, 0.25, _PPT_WHITE)
            _ppt_txt(slide, str(b)[:110], 1.8, by, W - 3.5, 1.0, 12, color=_PPT_GRAY_MED)
            by += 1.25


def _skills_cards(skills: list) -> list:
    if not skills:
        return [{"title": "기술 스택", "body": "내용을 입력하세요"}]
    cats  = ["Frontend / Backend", "DevOps / Infra", "기타 기술"]
    chunk = max(1, (len(skills) + 2) // 3)
    result = []
    for i in range(3):
        batch = skills[i * chunk:(i + 1) * chunk]
        if not batch:
            break
        result.append({"title": cats[i], "body": ", ".join(str(s) for s in batch)})
    return result or [{"title": "기술 스택", "body": "내용을 입력하세요"}]


# ── PPT Preview ───────────────────────────────────────────────────────────────

@app.post("/generate/ppt/preview")
def generate_ppt_preview(req: PptPreviewReq):
    sess = session_get(req.session_id)
    if sess is None:
        raise HTTPException(404, "세션을 찾을 수 없습니다.")

    all_projects: list = []
    all_skills:   list = []
    summary = ""
    for m in sess["materials"]:
        p = m["parsed"]
        all_projects.extend(p.get("projects", []))
        all_skills.extend(p.get("skills", []))
        if not summary and p.get("summary"):
            summary = p["summary"]

    ans_bullets = [str(a)[:100] for a in sess.get("answers", []) if a][:3]

    retro = sess.get("retro", {})
    retro_bullets: list = []
    if retro and retro.get("type") not in (None, "skip", ""):
        for pr in (retro.get("data") or []):
            if isinstance(pr, dict):
                for k, v in pr.items():
                    if k != "project" and v:
                        retro_bullets.append(f"[{k}] {str(v)[:80]}")
        retro_bullets = retro_bullets[:3]

    pt        = req.portfolio_type
    name_line = summary[:40] if summary else "이름 / 직함"

    def _base(sid, layout, title="", subtitle="", content=""):
        return {
            "id": sid, "layout": layout, "title": title,
            "subtitle": subtitle, "content": content,
            "bullets": [], "cards": [],
            "role": "", "description": "", "tech_stack": [], "achievements": [],
        }

    def _proj(idx, sid):
        if idx >= len(all_projects):
            return None
        p = all_projects[idx]
        return {
            "id": sid, "layout": "project",
            "title": p.get("name", f"프로젝트 {idx + 1}"),
            "subtitle": "", "content": "", "bullets": [], "cards": [],
            "role":         p.get("role", ""),
            "description":  p.get("description", ""),
            "tech_stack":   p.get("tech_stack", []),
            "achievements": p.get("achievements", []),
        }

    # (subtitle, s2_title, s2_content, s3_title, s4_title, s7_title, s8_title, cards_fn)
    CFG = {
        "developer": (
            "DEVELOPER PORTFOLIO", "기술 역량", "보유 기술 스택 및 개발 경험",
            "기술 스택", "주요 프로젝트", "기술적 문제해결", "회고 & 성장",
            lambda: _skills_cards(all_skills),
        ),
        "planner": (
            "PLANNER PORTFOLIO", "기획 역량", "서비스 기획 및 리서치 경험",
            "기획 방법론", "주요 프로젝트", "서비스 성과 지표", "회고 & 인사이트",
            lambda: [
                {"title": "리서치",    "body": "사용자 인터뷰, 설문, 경쟁사 분석"},
                {"title": "UX 설계",  "body": "정보구조, 와이어프레임, 프로토타입"},
                {"title": "지표 설계", "body": "MAU, 전환율, NPS, 리텐션"},
            ],
        ),
        "designer": (
            "DESIGNER PORTFOLIO", "디자인 역량", "툴 역량 및 UX 개선 경험",
            "툴 & 결과물", "주요 프로젝트", "디자인 프로세스", "회고 & 성장",
            lambda: [
                {"title": "디자인 툴", "body": "Figma, Sketch, Adobe XD, Illustrator"},
                {"title": "브랜딩",   "body": "BI 설계, 로고, 컬러 시스템, 타이포"},
                {"title": "UX 개선",  "body": "사용성 테스트, IA, 인터랙션 설계"},
            ],
        ),
        "marketer": (
            "MARKETER PORTFOLIO", "채널 & 캠페인 역량", "채널별 퍼포먼스 및 캠페인 경험",
            "채널별 경험", "주요 프로젝트", "핵심 퍼포먼스 지표", "회고 & 전략",
            lambda: [
                {"title": "SNS / 콘텐츠", "body": "Instagram, YouTube, 블로그, 바이럴"},
                {"title": "퍼포먼스",     "body": "Google Ads, Meta Ads, ROAS, CTR"},
                {"title": "데이터 분석",  "body": "GA4, Amplitude, A/B 테스트, CAC"},
            ],
        ),
    }
    sub, s2t, s2c, s3t, s4t, s7t, s8t, cards_fn = CFG.get(pt, CFG["developer"])

    slides = [
        {**_base("s1", "title",   name_line, sub,  summary[:120])},
        {**_base("s2", "section", s2t,       content=s2c)},
        {**_base("s3", "content", s3t), "cards": cards_fn()},
        {**_base("s4", "section", s4t)},
    ]
    for idx, sid in ((0, "s5"), (1, "s6")):
        sl = _proj(idx, sid)
        if sl:
            slides.append(sl)
    slides += [
        {**_base("s7", "content", s7t), "bullets": ans_bullets or ["답변 내용을 입력하세요"]},
        {**_base("s8", "content", s8t), "bullets": retro_bullets},
        {**_base("s9", "closing", "감사합니다", content=summary[:80] if summary else "")},
    ]
    return {"slides": slides, "portfolio_type": pt, "session_id": req.session_id}


# ── PPT Download ──────────────────────────────────────────────────────────────

@app.post("/generate/ppt/download")
def generate_ppt_download(req: PptDownloadReq):
    if not PPTX_OK:
        raise HTTPException(400, "python-pptx가 설치되지 않았습니다.")

    accent = _PPT_ACCENTS.get(req.portfolio_type, "6366F1")
    prs = Presentation()
    prs.slide_width  = Cm(33.87)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    for raw in req.slides:
        sd     = raw if isinstance(raw, dict) else dict(raw)
        layout = sd.get("layout", "content")
        slide  = prs.slides.add_slide(blank)
        if layout in ("title", "section", "closing"):
            _ppt_dark_slide(slide, sd, layout, accent)
        elif layout == "project":
            _ppt_project_slide(slide, sd, accent)
        else:
            _ppt_content_slide(slide, sd, accent)

    tmp = os.path.join(tempfile.gettempdir(), f"portfolio_{uuid.uuid4().hex}.pptx")
    prs.save(tmp)
    return FileResponse(
        tmp,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="portfolio.pptx",
        headers={"Content-Disposition": "attachment; filename=portfolio.pptx"},
        background=BackgroundTask(os.unlink, tmp),
    )


# ── Payment ───────────────────────────────────────────────────────────────────

@app.post("/payment/create-order")
def payment_create_order(req: PaymentCreateReq, user: dict = Depends(get_current_user)):
    order_id = f"portfolog-{uuid.uuid4().hex[:16]}"
    payment_create(user["sub"], order_id, req.amount)
    return {
        "order_id": order_id,
        "amount": req.amount,
        "order_name": "Portfolog Pro 1개월",
        "customer_name": user.get("email", "고객"),
    }


@app.post("/payment/confirm")
def payment_confirm_route(req: PaymentConfirmReq, user: dict = Depends(get_current_user)):
    status = payment_get_status(req.order_id)
    if status == "confirmed":
        return {"ok": True, "message": "이미 처리된 결제입니다."}

    try:
        toss_confirm(req.payment_key, req.order_id, req.amount)
    except ValueError as e:
        payment_fail(req.order_id)
        raise HTTPException(status_code=400, detail=str(e))

    payment_confirm(req.order_id, req.payment_key)
    profile_upgrade_to_pro(user["sub"])
    return {"ok": True, "message": "Pro 플랜으로 업그레이드 되었습니다!"}


@app.get("/payment/success")
def payment_success_page():
    return FileResponse("static/payment-success.html")


@app.get("/payment/fail")
def payment_fail_page():
    return FileResponse("static/payment-fail.html")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.getenv("PORT", 8000)))
